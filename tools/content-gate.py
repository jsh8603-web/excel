#!/usr/bin/env python3
"""
content-gate.py — 덱의 숫자가 'excel' 레포의 정답(tidy)으로 추적되는지 검증.

설계: rendering 게이트는 '보이는 게 맞나'를 본다. 이 게이트는 '적힌 게 맞나'를 본다.
무거운 절반(숫자 정규화·단위 스케일·괄호음수/△/%)은 jsh8603-web/excel 의 fpna.ingest.normalize 가
이미 구현했으므로 **재구현하지 않고 import 한다**. 이 파일은 얇은 크로스-레포 어댑터일 뿐이다.

흐름:
  1) excel 레포: py main.py ingest <원본.xlsx>  →  tidy.csv (정답 숫자 집합)
  2) 덱(HTML/PPTX)에서 숫자 토큰 추출 → 같은 normalize 로 정규화
  3) 각 덱 숫자가 tidy 정답 집합(native 또는 base-scaled)에 허용오차 내로 존재하는지 확인
  4) 어디에도 없는 숫자 = '미검증(uncited/fabricated/오타)' → 보고 + exit 1

휴리스틱(정직): 슬라이드 숫자가 어느 metric을 가리키는지까지 의미매칭하지 않는다.
'소스 집합 멤버십'만 본다 — 자릿수 전치/허위 수치는 어떤 정답과도 안 맞아 잡히고,
정상 인용은 통과한다. 의미적 metric 정합은 후속 확장.
"""
from __future__ import annotations
import csv, os, re, sys, argparse, html as _html

# --- excel 레포의 검증된 파서를 재사용 (재구현 금지) -------------------------------
def _load_excel_normalize(excel_repo: str):
    sys.path.insert(0, excel_repo)
    from fpna.ingest.normalize import normalize_value, scale_for_unit, split_cell_scale  # type: ignore
    return normalize_value, scale_for_unit, split_cell_scale

# 숫자 토큰: (1,234) △30 ▲5 -50 1,200 85% 3.4억 1,200천원 (50백만) ...
_UNITS = ["천원", "백만원", "백만", "억원", "억", "조", "만원", "만", "천", "원", "%"]
_NUM_TOKEN = re.compile(
    r"[(△▲\-+]?\s*[\d][\d,]*\.?\d*\s*\)?\s*(?:%s)?\s*\)?" % "|".join(sorted(_UNITS, key=len, reverse=True))
)

def deck_text_from_html(path: str) -> str:
    raw = open(path, encoding="utf-8").read()
    raw = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=re.S | re.I)
    return _html.unescape(re.sub(r"<[^>]+>", " ", raw))

def deck_text_from_pptx(path: str) -> str:
    from pptx import Presentation  # python-pptx
    prs = Presentation(path)
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
    return "\n".join(out)

def deck_text(path: str) -> str:
    return deck_text_from_pptx(path) if path.lower().endswith(".pptx") else deck_text_from_html(path)


def extract_numbers(text: str, normalize_value, scale_for_unit, split_cell_scale):
    """덱 텍스트 → [(원토큰, magnitude, had_unit)].
    셀내 단위(억/조/천원/백만)는 excel 의 split_cell_scale 로 분해(견고)."""
    out = []
    for m in _NUM_TOKEN.finditer(text):
        tok = m.group(0).strip()
        if not re.search(r"\d", tok):
            continue
        body, scale = split_cell_scale(tok)       # ('3.4', 1e8) / ('(50', 1e6) / ('1,200', 1)
        if body.count("(") > body.count(")"):       # 단위와 함께 닫는 괄호가 떨어진 경우 보정
            body += ")"
        bare = body.strip().lstrip("(△▲-+ ").rstrip(") ")
        # 연도/기간 마커 스킵: 단위 없는 4자리 정수(1900~2100)
        if scale == 1 and re.fullmatch(r"\d{4}", bare) and 1900 <= int(bare) <= 2100:
            continue
        val = normalize_value(body)
        num = val[0] if isinstance(val, tuple) else val
        if num is None:
            continue
        had_unit = scale != 1
        out.append((tok, abs(float(num) * (scale if had_unit else 1)), had_unit))
    return out

def tidy_truth_sets(tidy_csv: str, scale_for_unit):
    """tidy.csv → (native 절대값 집합, base-scaled 절대값 집합)."""
    native, base = set(), set()
    with open(tidy_csv, encoding="utf-8-sig", newline="") as f:
        for row in csv.DictReader(f):
            v = (row.get("value") or "").strip()
            if v == "":
                continue
            try:
                num = float(v)
            except ValueError:
                continue
            unit = (row.get("unit") or "").strip()
            scale = scale_for_unit(unit) if unit else 1
            native.add(round(abs(num), 6))
            base.add(round(abs(num) * scale, 6))
    return native, base

def matches(value: float, truth: set, rel_tol=0.001, abs_floor=0.5) -> bool:
    for t in truth:
        tol = max(abs(t) * rel_tol, abs_floor)
        if abs(value - t) <= tol:
            return True
    return False

def check_deck(deck_path: str, tidy_csv: str, excel_repo: str, ignore_small=10.0):
    normalize_value, scale_for_unit, split_cell_scale = _load_excel_normalize(excel_repo)
    native, base = tidy_truth_sets(tidy_csv, scale_for_unit)
    nums = extract_numbers(deck_text(deck_path), normalize_value, scale_for_unit, split_cell_scale)
    unverified = []
    for tok, val, had_unit in nums:
        if val < ignore_small:        # 페이지번호·축눈금 등 노이즈
            continue
        truth = base if had_unit else native   # 단위 있으면 base, 없으면 표시값(native)
        if not matches(val, truth):
            unverified.append((tok, val))
    return {"checked": len(nums), "unverified": unverified,
            "truth_size": len(native) + len(base)}

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck", help="슬라이드 HTML 또는 PPTX")
    ap.add_argument("tidy", help="excel 레포 ingest 산출 tidy.csv")
    ap.add_argument("--excel-repo", default=os.environ.get("EXCEL_REPO", "../excel"))
    ap.add_argument("--ignore-small", type=float, default=10.0)
    a = ap.parse_args()
    r = check_deck(a.deck, a.tidy, a.excel_repo, a.ignore_small)
    print(f"checked={r['checked']} truth_set={r['truth_size']} unverified={len(r['unverified'])}")
    for tok, val in r["unverified"]:
        print(f"  ✗ 미검증 숫자: '{tok}' (={val:g}) — 어떤 출처값과도 불일치")
    if r["unverified"]:
        print("CONTENT GATE FAIL: 출처로 추적 안 되는 숫자가 있음")
        sys.exit(1)
    print("CONTENT GATE PASS: 모든 유효 숫자가 출처로 추적됨")

if __name__ == "__main__":
    main()
